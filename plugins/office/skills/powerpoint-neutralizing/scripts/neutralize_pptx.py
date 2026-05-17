#!/usr/bin/env python3
"""Neutralize PPTX visual styling: uniform background, fonts, colors, no shadows."""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

# ── Defaults — override via CLI flags ──────────────────────────────────────────
BG_COLOR = "#cccccc"
FONT_COLOR = "#000000"
TITLE_FONT = "Helvetica"
TITLE_SIZE = 18
BODY_FONT = "Helvetica"
BODY_SIZE = 14
# ──────────────────────────────────────────────────────────────────────────────


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def remove_shadow(run) -> None:
    rPr = run._r.get_or_add_rPr()
    for tag in (qn("a:effectLst"), qn("a:effectDag")):
        el = rPr.find(tag)
        if el is not None:
            rPr.remove(el)


def is_title_shape(shape) -> bool:
    try:
        return shape.placeholder_format is not None and shape.placeholder_format.idx == 0
    except Exception:
        return False


def neutralize(
    input_path: str,
    output_path: str,
    bg_color: str,
    font_color: str,
    title_font: str,
    title_size: int,
    body_font: str,
    body_size: int,
) -> None:
    prs = Presentation(input_path)
    bg_rgb = hex_to_rgb(bg_color)
    text_rgb = hex_to_rgb(font_color)

    for slide in prs.slides:
        set_slide_background(slide, bg_rgb)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            title = is_title_shape(shape)
            font_name = title_font if title else body_font
            font_pt = Pt(title_size if title else body_size)

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = font_name
                    run.font.size = font_pt
                    run.font.color.rgb = text_rgb
                    remove_shadow(run)

    prs.save(output_path)
    print(f"Saved: {output_path}")


def main() -> None:
    p = argparse.ArgumentParser(description="Neutralize PPTX styling")
    p.add_argument("input", help="Input .pptx file")
    p.add_argument("-o", "--output", help="Output path (default: <input>-neutralized.pptx)")
    p.add_argument("--bg-color", default=BG_COLOR, metavar="HEX", help=f"Slide background color (default: {BG_COLOR})")
    p.add_argument("--font-color", default=FONT_COLOR, metavar="HEX", help=f"Text color (default: {FONT_COLOR})")
    p.add_argument("--title-font", default=TITLE_FONT, help=f"Title font name (default: {TITLE_FONT})")
    p.add_argument("--title-size", type=int, default=TITLE_SIZE, help=f"Title font size in pt (default: {TITLE_SIZE})")
    p.add_argument("--body-font", default=BODY_FONT, help=f"Body font name (default: {BODY_FONT})")
    p.add_argument("--body-size", type=int, default=BODY_SIZE, help=f"Body font size in pt (default: {BODY_SIZE})")
    args = p.parse_args()

    stem = Path(args.input).stem
    output = args.output or str(Path(args.input).with_stem(stem + "-neutralized"))

    neutralize(
        input_path=args.input,
        output_path=output,
        bg_color=args.bg_color,
        font_color=args.font_color,
        title_font=args.title_font,
        title_size=args.title_size,
        body_font=args.body_font,
        body_size=args.body_size,
    )


if __name__ == "__main__":
    main()
